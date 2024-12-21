import os
from pptx import Presentation
from gtts import gTTS
from PIL import Image
import tempfile
import logging
import subprocess
from typing import List, Dict

logger = logging.getLogger(__name__)

class PPTProcessor:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
        logger.info(f"Initialized PPTProcessor with temp directory: {self.temp_dir}")
        
    def extract_text_from_slides(self, ppt_path):
        prs = Presentation(ppt_path)
        slides_text = []
        
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            slides_text.append(text.strip())
            
        return slides_text
    
    def create_audio_for_script(self, script_texts: List[str]) -> List[str]:
        """Create audio files from edited script texts"""
        audio_files = []
        for i, text in enumerate(script_texts):
            if text:
                try:
                    audio_path = os.path.join(self.temp_dir, f'slide_{i}.mp3')
                    tts = gTTS(text=text, lang='en')
                    tts.save(audio_path)
                    audio_files.append(audio_path)
                except Exception as e:
                    logger.error(f"Error creating audio for slide {i}: {str(e)}")
                    raise
            else:
                audio_files.append(None)
        return audio_files

    def extract_slides_as_images(self, ppt_path):
        image_files = []
        try:
            prs = Presentation(ppt_path)
            
            # Calculate slide size (assuming 96 DPI)
            width = int(prs.slide_width * 96 / 914400)  # convert EMU to pixels
            height = int(prs.slide_height * 96 / 914400)
            
            for i, slide in enumerate(prs.slides):
                # Create a blank image with white background
                img = Image.new('RGB', (width, height), 'white')
                
                # Process shapes in the slide (basic text rendering)
                y_offset = 50  # Start text from top with some margin
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Draw text on image
                        from PIL import ImageDraw
                        draw = ImageDraw.Draw(img)
                        text = shape.text.strip()
                        draw.text((50, y_offset), text, fill='black')
                        y_offset += 30  # Move down for next text block
                
                # Save the image
                img_path = os.path.join(self.temp_dir, f'slide_{i}.png')
                img.save(img_path, 'PNG')
                image_files.append(img_path)
                
            logger.info(f"Successfully extracted {len(image_files)} slides as images")
            return image_files
            
        except Exception as e:
            logger.error(f"Error extracting slides as images: {str(e)}")
            raise

    def combine_slide_and_audio(self, image_path: str, audio_path: str, output_path: str, duration: int = None):
        """Combine single image and audio using ffmpeg"""
        if not duration:
            # Get audio duration using ffprobe
            cmd = ['ffprobe', '-i', audio_path, '-show_entries', 'format=duration', '-v', 'quiet', '-of', 'csv="p=0"']
            duration = float(subprocess.check_output(cmd).decode().strip())

        # Combine image and audio using ffmpeg
        cmd = [
            'ffmpeg', '-y',
            '-loop', '1',
            '-i', image_path,
            '-i', audio_path,
            '-c:v', 'libx264',
            '-tune', 'stillimage',
            '-c:a', 'aac',
            '-b:a', '192k',
            '-pix_fmt', 'yuv420p',
            '-t', str(duration),
            output_path
        ]
        subprocess.run(cmd, check=True)
        return output_path

    def process_ppt(self, ppt_path: str, output_path: str, edited_script: List[str] = None) -> str:
        """Process PPT with optional edited script"""
        # Extract slides as images
        image_files = self.extract_slides_as_images(ppt_path)
        
        # Use edited script if provided, otherwise extract from slides
        script_texts = edited_script if edited_script else self.extract_text_from_slides(ppt_path)
        audio_files = self.create_audio_for_script(script_texts)

        # Create temporary video segments
        video_segments = []
        for i, (img_path, audio_path) in enumerate(zip(image_files, audio_files)):
            if audio_path:
                segment_path = os.path.join(self.temp_dir, f'segment_{i}.mp4')
                self.combine_slide_and_audio(img_path, audio_path, segment_path)
                video_segments.append(segment_path)

        # Concatenate all segments
        with open(os.path.join(self.temp_dir, 'segments.txt'), 'w') as f:
            for segment in video_segments:
                f.write(f"file '{segment}'\n")

        # Final concatenation
        cmd = [
            'ffmpeg', '-y',
            '-f', 'concat',
            '-safe', '0',
            '-i', os.path.join(self.temp_dir, 'segments.txt'),
            '-c', 'copy',
            output_path
        ]
        subprocess.run(cmd, check=True)

        # Cleanup
        for file in os.listdir(self.temp_dir):
            os.remove(os.path.join(self.temp_dir, file))
        os.rmdir(self.temp_dir)

        return output_path
