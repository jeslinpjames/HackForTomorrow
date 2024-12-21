import os
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# Load environment variables
load_dotenv(os.path.join(os.path.dirname(os.path.dirname(__file__)), '.env'))

UPLOAD_FOLDER = 'uploads'
AUDIO_FOLDER = 'static/audios'
IMAGE_FOLDER = 'static/images'  # New folder for slide images

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(AUDIO_FOLDER, exist_ok=True)
os.makedirs(IMAGE_FOLDER, exist_ok=True)

# Configure Gemini with API key from environment variable
api_key = os.getenv('GEMINI_API_KEY')
if not api_key:
    raise ValueError("GEMINI_API_KEY not found in environment variables")
genai.configure(api_key=api_key)

def slide_to_image(slide, image_path):
    """Convert a slide to an image using PIL with improved text rendering"""
    # Create a blank white image
    width = 1920  # 16:9 aspect ratio
    height = 1080
    image = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(image)
    
    # Try to load a default system font, fallback to default
    try:
        # For Windows
        font_large = ImageFont.truetype("arial.ttf", 60)
        font_normal = ImageFont.truetype("arial.ttf", 40)
    except:
        # Fallback to default
        font_large = ImageFont.load_default()
        font_normal = ImageFont.load_default()

    # Collect and organize text content
    title_text = ""
    body_texts = []
    
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
            
        # Check if shape is likely a title (usually first text box or at the top)
        if shape.top < height * 0.2 and not title_text:  # If in top 20% of slide
            title_text = shape.text.strip()
        else:
            if shape.text.strip():
                body_texts.append(shape.text.strip())
    
    # Draw title
    if title_text:
        # Center the title
        title_bbox = draw.textbbox((0, 0), title_text, font=font_large)
        title_width = title_bbox[2] - title_bbox[0]
        title_x = (width - title_width) // 2
        draw.text((title_x, 50), title_text, font=font_large, fill='black')
    
    # Draw body text
    y_position = 200  # Start body text below title
    for text in body_texts:
        # Wrap text if too long
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            current_line.append(word)
            test_line = ' '.join(current_line)
            bbox = draw.textbbox((0, 0), test_line, font=font_normal)
            if bbox[2] - bbox[0] > width - 100:  # Leave 50px margin on each side
                lines.append(' '.join(current_line[:-1]))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        # Draw each line
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=font_normal)
            line_width = bbox[2] - bbox[0]
            x_position = (width - line_width) // 2  # Center the line
            draw.text((x_position, y_position), line, font=font_normal, fill='black')
            y_position += 60  # Space between lines
        
        y_position += 40  # Extra space between paragraphs
    
    # Save the image
    image.save(image_path, 'PNG')

def generate_scripts(ppt_path):
    try:
        prs = Presentation(ppt_path)
        slides_data = []
        
        # Initialize Gemini model
        model = genai.GenerativeModel('gemini-1.5-flash')
        
        for index, slide in enumerate(prs.slides):
            # Save slide as image
            image_filename = f"slide_{index + 1}.png"
            image_path = os.path.join(IMAGE_FOLDER, image_filename)
            
            # Convert slide to image
            slide_to_image(slide, image_path)
            
            # Extract slide content
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            
            # Generate scene description using slide text and Gemini
            prompt = (
                f"Based on the following slide content: {slide_text}\n"
                "Describe this slide as if narrating to a blind user for 15 seconds. "
                "Provide a short and concise description. Focus on the key points and main message."
            )
            
            response = model.generate_content(prompt)
            scene_description = response.text.strip()
            
            # Generate voiceover
            audio_filename = f"slide_{index + 1}.mp3"
            audio_path = os.path.join(AUDIO_FOLDER, audio_filename)
            generate_voiceover(scene_description, audio_path)
            
            slides_data.append({
                'filename': os.path.basename(ppt_path),
                'slide_no': index + 1,
                'image_url': f"/static/images/{image_filename}",
                'script': scene_description,
                'audio_url': f"/static/audios/{audio_filename}"
            })
            
        return slides_data
    except Exception as e:
        raise Exception(f"Error processing PPT: {str(e)}")

def generate_voiceover(text, audio_path):
    # Implement text-to-speech conversion using a suitable library or API
    # Example using gTTS:
    from gtts import gTTS
    tts = gTTS(text)
    tts.save(audio_path)
